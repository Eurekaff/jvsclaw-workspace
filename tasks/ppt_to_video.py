#!/usr/bin/env python3
"""
PPT 动画转视频

流程：
1. 用 python-pptx 创建 PPT
2. 添加动画效果
3. 导出为图片
4. 合成视频
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image
import imageio.v2 as imageio
from pathlib import Path
from datetime import datetime
import io


class PPTVideoGenerator:
    """PPT 转视频"""
    
    def __init__(self, workspace_root: str):
        self.workspace_root = Path(workspace_root)
        self.output_dir = self.workspace_root / "artifacts" / "douyin_videos"
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.temp_dir = self.output_dir / "temp_ppt"
        self.temp_dir.mkdir(parents=True, exist_ok=True)
    
    def create_ppt(self):
        """创建 PPT"""
        
        print("\n📊 创建 PPT...")
        
        prs = Presentation()
        # 设置 9:16 比例（抖音竖屏）
        prs.slide_width = Inches(6.75)
        prs.slide_height = Inches(12)
        
        # 场景内容
        scenes = [
            ("⚠️别再手动工作了！", 3),
            ("每天加班到深夜？", 7),
            ("工具 1：自动写文档", 10),
            ("3 秒生成！", 5),
            ("工具 2：自动做 PPT", 10),
            ("效率翻倍！", 5),
            ("之前 3 小时", 5),
            ("现在 3 分钟", 5),
            ("评论区领工具👇", 10),
        ]
        
        for text, duration in scenes:
            # 创建幻灯片
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式
            
            # 添加文字
            textbox = slide.shapes.add_textbox(
                Inches(0.5), Inches(4),
                Inches(5.75), Inches(4)
            )
            tf = textbox.text_frame
            p = tf.paragraphs[0]
            p.text = text
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(48)
            p.font.bold = True
            
            print(f"  ✅ 场景：{text} ({duration}秒)")
        
        # 保存 PPT
        ppt_file = self.temp_dir / "video.pptx"
        prs.save(str(ppt_file))
        print(f"\n📁 PPT 已保存：{ppt_file}")
        
        return str(ppt_file)
    
    def export_slides(self, ppt_file: str):
        """导出幻灯片为图片"""
        
        print("\n📸 导出幻灯片...")
        
        # 注意：python-pptx 不能直接导出图片
        # 这里用简化方式：创建纯色图片代替
        from PIL import Image, ImageDraw, ImageFont
        
        scenes = [
            ("⚠️别再手动工作了！", 3),
            ("每天加班到深夜？", 7),
            ("工具 1：自动写文档", 10),
            ("3 秒生成！", 5),
            ("工具 2：自动做 PPT", 10),
            ("效率翻倍！", 5),
            ("之前 3 小时", 5),
            ("现在 3 分钟", 5),
            ("评论区领工具👇", 10),
        ]
        
        images = []
        
        for text, duration in scenes:
            # 创建图片（每张 2 秒，2fps = 4 张）
            frames_count = duration * 2
            
            for i in range(frames_count):
                img = Image.new('RGB', (1080, 1920), color=(26, 26, 46))
                draw = ImageDraw.Draw(img)
                
                # 绘制文字
                try:
                    font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 72)
                except:
                    font = ImageFont.load_default()
                
                # 计算文字位置
                bbox = draw.textbbox((0, 0), text, font=font)
                text_width = bbox[2] - bbox[0]
                text_height = bbox[3] - bbox[1]
                
                x = (1080 - text_width) // 2
                y = (1920 - text_height) // 2
                
                draw.text((x, y), text, font=font, fill=(255, 255, 255))
                
                images.append(img)
        
        print(f"  ✅ 生成 {len(images)} 帧")
        return images
    
    def create_video(self, images: list, fps: int = 2):
        """创建视频"""
        
        print("\n🎬 生成视频...")
        
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_file = self.output_dir / f"ppt_video_{timestamp}.mp4"
        
        # 转换 PIL 图像为 numpy
        import numpy as np
        frames = [np.array(img) for img in images]
        
        # 生成视频
        imageio.mimwrite(str(output_file), frames, fps=fps)
        
        size_mb = output_file.stat().st_size / 1024 / 1024
        
        print(f"\n{'='*60}")
        print("✅ 视频生成完成！")
        print(f"{'='*60}")
        print(f"📁 文件：{output_file}")
        print(f"⏱️  时长：{len(frames)/fps:.1f}秒")
        print(f"📊 大小：{size_mb:.2f} MB")
        
        return {
            "file": str(output_file),
            "duration": len(frames) / fps,
            "size_mb": size_mb
        }
    
    def generate(self):
        """完整流程"""
        
        print("="*60)
        print("🎬 PPT 动画转视频")
        print("="*60)
        
        # 创建 PPT
        ppt_file = self.create_ppt()
        
        # 导出图片
        images = self.export_slides(ppt_file)
        
        # 创建视频
        result = self.create_video(images)
        
        return result


if __name__ == "__main__":
    generator = PPTVideoGenerator("/home/admin/openclaw/workspace")
    result = generator.generate()
    
    if result:
        print(f"\n🎉 成功：{result['file']}")
